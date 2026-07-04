#!/usr/bin/env python3
"""Build the metaDEBASS hackathon results deck.

Run:
    source ~/.venvs/debass_py313/bin/activate
    python presentations/hackathon_20260612/build_deck.py

Output:
    presentations/hackathon_20260612/metaDEBASS_hackathon_2026-06-12.pptx

All science metrics are read from facts.json. Figures are embedded from figs/.
"""
from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parent
FACTS_PATH = BASE / "facts.json"
FIG_DIR = BASE / "figs"
OUT_PATH = BASE / "metaDEBASS_hackathon_2026-06-12.pptx"


# Palette mirrors presentations/build_v7_status.py.
BLACK = RGBColor(0x00, 0x00, 0x00)
NAVY = RGBColor(0x10, 0x2A, 0x43)
SLATE = RGBColor(0x33, 0x44, 0x55)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
SOFT = RGBColor(0xCB, 0xD5, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x1D, 0x4E, 0xD8)
GREEN = RGBColor(0x05, 0x80, 0x4F)
RED = RGBColor(0xB9, 0x1C, 0x1C)
AMBER = RGBColor(0xB4, 0x53, 0x09)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
RED_LIGHT = RGBColor(0xFE, 0xE2, 0xE2)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def load_facts() -> dict:
    with FACTS_PATH.open("r", encoding="utf-8") as f:
        return json.load(f)


facts = load_facts()


def num(value: float, *, plus: bool = False, digits: int = 3) -> str:
    text = f"{value:.{digits}f}"
    if plus and value > 0:
        text = f"+{text}"
    return text


def whole(value: int) -> str:
    return f"{value:,}"


def ci_text(ci: list[float] | tuple[float, float] | None) -> str:
    if ci is None:
        return "calibration / guardrail"
    return f"[{num(ci[0])}, {num(ci[1])}]"


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=14,
    bold=False,
    color=BLACK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    italic=False,
    font="Helvetica",
    margin=0.03,
):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(0)
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.font.name = font
    return tb


def add_rect(slide, x, y, w, h, fill=LIGHT, line=None, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_rule(slide, x, y, w, color=SLATE, height=0.018):
    return add_rect(slide, x, y, w, height, fill=color)


def add_arrow(slide, x1, y1, x2, y2, color=GREY, width=1.4):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    line_fmt = conn.line._get_or_add_ln()
    from lxml import etree
    from pptx.oxml.ns import qn

    tail = etree.SubElement(line_fmt, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return conn


def add_kicker(slide, text, x=0.5, y=0.26, w=12.3):
    add_rect(slide, x, y + 0.09, 0.28, 0.045, fill=ACCENT)
    add_textbox(
        slide,
        x + 0.38,
        y,
        w - 0.38,
        0.28,
        text.upper(),
        size=8.5,
        bold=True,
        color=GREY,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def slide_header(slide, kicker, title, subtitle=None):
    add_kicker(slide, kicker)
    add_textbox(slide, 0.5, 0.53, 12.25, 0.55, title, size=25, bold=True)
    if subtitle:
        add_textbox(slide, 0.5, 1.08, 12.25, 0.38, subtitle, size=12.5, color=GREY)
    add_rule(slide, 0.5, 1.52, 12.3, color=SLATE)


def add_table(
    slide,
    x,
    y,
    w,
    h,
    headers,
    rows,
    *,
    header_fill=NAVY,
    header_color=WHITE,
    font_size=10.5,
    header_size=10.5,
    col_widths=None,
    num_cols=None,
    first_col_bold=False,
):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    if col_widths is not None:
        total = sum(col_widths)
        for ci, frac in enumerate(col_widths):
            tbl.columns[ci].width = Inches(w * frac / total)
    num_cols = num_cols or set()
    for ci, header in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.font.size = Pt(header_size)
        p.font.bold = True
        p.font.color.rgb = header_color
        p.font.name = "Helvetica"
        p.alignment = PP_ALIGN.RIGHT if ci in num_cols else PP_ALIGN.LEFT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.07)
        cell.margin_right = Inches(0.07)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.text = str(val)
            p.font.size = Pt(font_size)
            p.font.color.rgb = BLACK
            p.font.name = "Helvetica"
            p.font.bold = first_col_bold and ci == 0
            p.alignment = PP_ALIGN.RIGHT if ci in num_cols else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
    return tbl


def add_metric(slide, x, y, w, h, value, label, context="", *, color=NAVY, fill=LIGHT):
    add_rect(slide, x, y, w, h, fill=fill, line=SOFT)
    add_textbox(slide, x + 0.16, y + 0.08, w - 0.32, 0.45, value, size=22, bold=True, color=color)
    add_textbox(slide, x + 0.16, y + 0.56, w - 0.32, 0.36, label, size=10.5, bold=True, color=BLACK)
    if context:
        add_textbox(slide, x + 0.16, y + 0.93, w - 0.32, h - 1.0, context, size=9.5, color=GREY)


def add_picture_fit(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as im:
        ratio = im.width / im.height
    box_ratio = w / h
    if ratio >= box_ratio:
        pic_w = w
        pic_h = w / ratio
        pic_x = x
        pic_y = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * ratio
        pic_x = x + (w - pic_w) / 2
        pic_y = y
    return slide.shapes.add_picture(str(path), Inches(pic_x), Inches(pic_y), Inches(pic_w), Inches(pic_h))


def footer(slide, text="metaDEBASS hackathon results"):
    add_textbox(slide, 0.5, 7.18, 12.3, 0.16, text, size=6.5, color=GREY, align=PP_ALIGN.CENTER)


def bullet_block(slide, x, y, w, h, title, body, *, fill=LIGHT, title_color=NAVY):
    add_rect(slide, x, y, w, h, fill=fill, line=SOFT)
    add_textbox(slide, x + 0.18, y + 0.14, w - 0.36, 0.34, title, size=13, bold=True, color=title_color)
    add_textbox(slide, x + 0.18, y + 0.58, w - 0.36, h - 0.7, body, size=10.6, color=BLACK)


def require_assets():
    required = [
        FIG_DIR / "fig_headline.png",
        FIG_DIR / "fig_auc_vs_ndet.png",
        FIG_DIR / "fig_dp1_ef.png",
        FIG_DIR / "fig_trust_heads.png",
    ]
    missing = [str(p) for p in required if not p.exists()]
    if missing:
        raise FileNotFoundError("Missing required figure(s): " + ", ".join(missing))


def add_title_slide():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
    add_rule(s, 0.5, 0.58, 12.3, color=NAVY, height=0.035)
    add_textbox(
        s,
        0.55,
        1.25,
        11.9,
        0.78,
        "metaDEBASS \u2014 who do you trust tonight?",
        size=32,
        bold=True,
        color=BLACK,
    )
    add_textbox(
        s,
        0.55,
        2.07,
        11.9,
        0.86,
        "Trust-aware fusion of transient classifiers for Rubin/LSST follow-up",
        size=23,
        color=SLATE,
    )
    add_textbox(
        s,
        0.55,
        3.05,
        11.7,
        0.38,
        "Hackathon results  |  June 12, 2026",
        size=14.5,
        color=GREY,
    )
    add_rect(s, 0.55, 4.05, 12.25, 1.68, fill=NAVY)
    add_textbox(
        s,
        0.85,
        4.25,
        11.65,
        0.38,
        "Headline on the locked spectroscopic test",
        size=13,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    h = facts["headline"]
    add_metric(
        s,
        0.9,
        4.78,
        3.5,
        0.78,
        num(h["v9c_delta_vs_v6e2"]["delta"], plus=True),
        "macro-AUC delta vs v6e2",
        f"{ci_text([h['v9c_delta_vs_v6e2']['lo'], h['v9c_delta_vs_v6e2']['hi']])}",
        color=GREEN,
        fill=WHITE,
    )
    add_metric(
        s,
        4.9,
        4.78,
        3.5,
        0.78,
        num(h["v9c_macro_auc"]["value"]),
        "v9c macro AUC",
        "object-level, spectroscopic-only",
        color=ACCENT,
        fill=WHITE,
    )
    add_metric(
        s,
        8.9,
        4.78,
        3.5,
        0.78,
        whole(h["n_test_objects"]),
        "locked test objects",
        "same split across versions",
        color=NAVY,
        fill=WHITE,
    )
    add_textbox(
        s,
        0.55,
        6.27,
        12.25,
        0.4,
        facts["headline"]["protocol"],
        size=10.5,
        color=GREY,
        align=PP_ALIGN.CENTER,
        italic=True,
    )


def add_problem_slide():
    s = prs.slides.add_slide(BLANK)
    slide_header(
        s,
        "Problem",
        "Early follow-up is a trust problem, not just a classifier problem.",
        "Observers must decide before the light curve is mature and while broker classifiers disagree.",
    )
    add_rect(s, 0.65, 1.9, 5.25, 4.85, fill=NAVY)
    add_textbox(s, 0.95, 2.18, 4.65, 0.42, "The operating pressure", size=16, bold=True, color=WHITE)
    add_textbox(
        s,
        0.95,
        2.82,
        4.65,
        3.35,
        facts["narrative"]["problem"],
        size=20,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        s,
        0.95,
        6.23,
        4.65,
        0.3,
        "Source: facts.narrative.problem",
        size=8,
        color=SOFT,
        italic=True,
    )
    rail_x = 6.35
    bullet_block(
        s,
        rail_x,
        1.9,
        6.35,
        1.35,
        "Millions of alerts",
        "The alert stream is large enough that triage quality matters immediately.",
        fill=LIGHT,
    )
    bullet_block(
        s,
        rail_x,
        3.48,
        6.35,
        1.35,
        "Many disagreeing experts",
        "The product needs to expose which expert is credible for this object and epoch.",
        fill=LIGHT,
    )
    bullet_block(
        s,
        rail_x,
        5.06,
        6.35,
        1.35,
        "Decisions at 3-5 detections",
        "The system must work while the evidence is still sparse.",
        fill=LIGHT,
    )
    footer(s)


def add_architecture_slide():
    s = prs.slides.add_slide(BLANK)
    slide_header(
        s,
        "What we built",
        "metaDEBASS is a meta-layer that turns expert disagreement into calibrated trust.",
        "The primary output is per-expert confidence; priority lists are a downstream use case.",
    )
    y = 2.08
    box_h = 1.25
    boxes = [
        (0.45, 2.05, "lightcurve +\n29 broker/local experts", BLUE_LIGHT),
        (2.95, 2.25, "pooled trust model\nP(expert right | object, epoch)", LIGHT),
        (5.65, 2.05, "multiclass fusion +\nDirichlet calibration", LIGHT),
        (8.35, 1.65, "conformal sets", LIGHT),
        (10.4, 2.0, "goal-conditioned\npriority lists\nIa / non-Ia / other", GREEN_LIGHT),
    ]
    for x, w, txt, fill in boxes:
        add_rect(s, x, y, w, box_h, fill=fill, line=SOFT)
        add_textbox(s, x + 0.12, y + 0.13, w - 0.24, box_h - 0.26, txt, size=11.8, bold=True, color=BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for x1, x2 in [(2.5, 2.95), (5.2, 5.65), (7.7, 8.35), (10.0, 10.4)]:
        add_arrow(s, x1, y + box_h / 2, x2, y + box_h / 2)
    add_rect(s, 0.65, 4.45, 12.05, 1.55, fill=NAVY)
    add_textbox(s, 0.95, 4.67, 11.45, 1.05, facts["narrative"]["solution"], size=16.5, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 0.65, 6.42, 12.05, 0.34, "Architecture is deliberately goal-conditioned: SN Ia, non-Ia, and 'other' users can ask different follow-up questions.", size=11.5, color=SLATE, align=PP_ALIGN.CENTER)
    footer(s)


def add_honesty_slide():
    s = prs.slides.add_slide(BLANK)
    slide_header(
        s,
        "Honest-by-construction",
        "The deck reports both the headline and the things the gates rejected.",
        "The strongest result is useful only because the evaluation contract is explicit.",
    )
    add_rect(s, 0.55, 1.9, 12.25, 1.35, fill=NAVY)
    add_textbox(s, 0.85, 2.07, 11.65, 0.95, facts["narrative"]["honesty"], size=14.5, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bug = facts["scale"]["label_bug_fixed"]
    add_rect(s, 0.55, 3.72, 12.25, 1.05, fill=RED_LIGHT, line=RED)
    add_textbox(s, 0.85, 3.88, 2.25, 0.36, "Bug found + fixed", size=13, bold=True, color=RED)
    add_textbox(s, 3.05, 3.82, 9.35, 0.52, bug, size=14, bold=True, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
    items = [
        ("Pre-registered", "headline + guards"),
        ("Locked test", "byte-identical split"),
        ("Cal-decided", "component gates with bootstrap CIs"),
        ("Audited", "adversarial recomputation + code audit"),
    ]
    x = 0.55
    for label, body in items:
        bullet_block(s, x, 5.25, 2.86, 1.12, label, body, fill=LIGHT)
        x += 3.13
    footer(s)


def add_headline_slide():
    s = prs.slides.add_slide(BLANK)
    h = facts["headline"]
    slide_header(
        s,
        "Headline result",
        "v9c improves the locked spectroscopic headline by +0.129 macro AUC.",
        h["protocol"],
    )
    add_picture_fit(s, FIG_DIR / "fig_headline.png", 0.45, 1.82, 8.15, 4.65)
    add_metric(
        s,
        8.92,
        1.92,
        3.55,
        1.1,
        num(h["v9c_delta_vs_v6e2"]["delta"], plus=True),
        "v9c delta vs v6e2",
        ci_text([h["v9c_delta_vs_v6e2"]["lo"], h["v9c_delta_vs_v6e2"]["hi"]]),
        color=GREEN,
    )
    add_metric(
        s,
        8.92,
        3.26,
        3.55,
        1.1,
        num(h["v9c_macro_auc"]["value"]),
        "macro AUC",
        ci_text([h["v9c_macro_auc"]["lo"], h["v9c_macro_auc"]["hi"]]),
        color=ACCENT,
    )
    add_metric(
        s,
        8.92,
        4.6,
        3.55,
        1.1,
        whole(h["n_test_objects"]),
        "spectroscopic test objects",
        "locked object-level test split",
        color=NAVY,
    )
    add_metric(
        s,
        8.92,
        5.94,
        3.55,
        0.78,
        num(h["v9c_snia_auc"]),
        "SN Ia AUC",
        "n_det=5",
        color=AMBER,
    )
    footer(s)


def add_early_epoch_slide():
    s = prs.slides.add_slide(BLANK)
    early = facts["auc_vs_ndet_snia"]
    slide_header(
        s,
        "The gain is present at the 3-5 detection operating point.",
        "Spec-only SN Ia AUC stays high before the light curve is mature.",
        "Values below are read directly from facts.auc_vs_ndet_snia.",
    )
    add_picture_fit(s, FIG_DIR / "fig_auc_vs_ndet.png", 0.45, 1.78, 8.9, 4.9)
    f_ndet = early["fusion_v8"]["n_det"]
    v_ndet = early["v6e2_rescored"]["n_det"]
    rows = []
    for target in (3, 5):
        fi = f_ndet.index(target)
        vi = v_ndet.index(target)
        rows.append(
            (
                f"n_det={target}",
                num(early["fusion_v8"]["auc_snia"][fi]),
                num(early["v6e2_rescored"]["auc_snia"][vi]),
                num(early["fusion_v8"]["auc_snia"][fi] - early["v6e2_rescored"]["auc_snia"][vi], plus=True),
            )
        )
    add_table(
        s,
        9.65,
        2.03,
        3.1,
        1.65,
        ["Epoch", "v9c", "v6e2", "Delta"],
        rows,
        font_size=10.5,
        header_size=9.5,
        col_widths=[1.0, 0.7, 0.7, 0.7],
        num_cols={1, 2, 3},
    )
    add_rect(s, 9.65, 4.18, 3.1, 1.92, fill=NAVY)
    add_textbox(s, 9.88, 4.38, 2.64, 0.36, "Why it matters", size=13, bold=True, color=WHITE)
    add_textbox(
        s,
        9.88,
        4.86,
        2.64,
        0.92,
        "This is exactly when a follow-up observer has to choose tonight's targets.",
        size=12,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    footer(s)


def add_dp1_slide():
    s = prs.slides.add_slide(BLANK)
    dp1 = facts["dp1_enrichment_top1pct"]
    slide_header(
        s,
        "On real Rubin DP1 data, the ranker suppresses contaminants and recovers SNe.",
        "EF=1 is random; contaminant classes want low, Published SNe wants high.",
        dp1["note"],
    )
    add_picture_fit(s, FIG_DIR / "fig_dp1_ef.png", 0.35, 1.78, 8.95, 5.05)
    ecl_prev = dp1["v6e2"]["EclBin+RRLyrae"]
    ecl_new = dp1["fusion_v9c"]["EclBin+RRLyrae"]["ef"]
    gaia_prev = dp1["v6e2"]["Gaia variables"]
    gaia_new = dp1["fusion_v9c"]["Gaia variables"]["ef"]
    sne_prev = dp1["v6e2"]["Published SNe"]
    sne_new = dp1["fusion_v9c"]["Published SNe"]["ef"]
    add_metric(
        s,
        9.55,
        1.92,
        3.05,
        1.08,
        f"{num(ecl_prev)} -> {num(ecl_new)}",
        "EclBin+RRLyrae",
        "periodic-variable contamination moves down",
        color=GREEN,
    )
    add_metric(
        s,
        9.55,
        3.26,
        3.05,
        1.08,
        f"{num(gaia_prev)} -> {num(gaia_new)}",
        "Gaia variables",
        "another contaminant class moves down",
        color=GREEN,
    )
    add_metric(
        s,
        9.55,
        4.6,
        3.05,
        1.08,
        f"{num(sne_prev)} -> {num(sne_new)}",
        "Published SNe",
        "the ranker finds known SNe",
        color=ACCENT,
    )
    footer(s)


def add_trust_slide():
    s = prs.slides.add_slide(BLANK)
    seq = next(row for row in facts["trust_heads"] if row["expert"] == "seq_v9")
    scale = facts["scale"]
    slide_header(
        s,
        "The trust deliverable is calibrated per-expert confidence.",
        "Each bar is a trust head: P(expert right | object, epoch), not a fused class posterior.",
        "seq_v9 is highlighted because it is the broadest local sequence expert.",
    )
    add_picture_fit(s, FIG_DIR / "fig_trust_heads.png", 0.35, 1.78, 8.75, 5.05)
    add_metric(
        s,
        9.4,
        1.92,
        3.2,
        1.05,
        whole(scale["experts_trust_headed"]),
        "trust-headed experts",
        "from facts.scale.experts_trust_headed",
        color=NAVY,
    )
    add_metric(
        s,
        9.4,
        3.22,
        3.2,
        1.05,
        num(seq["cal_auc"]),
        "seq_v9 trust AUC",
        f"n_test={whole(seq['n_test'])}",
        color=RED,
    )
    add_rect(s, 9.4, 4.72, 3.2, 1.42, fill=NAVY)
    add_textbox(s, 9.62, 4.92, 2.76, 0.34, "Product contract", size=13, bold=True, color=WHITE)
    add_textbox(
        s,
        9.62,
        5.36,
        2.76,
        0.54,
        "The emitted payload is expert_confidence per expert, epoch, and object.",
        size=10.8,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    footer(s)


def add_sequence_expert_slide():
    s = prs.slides.add_slide(BLANK)
    scale = facts["scale"]
    seq = next(row for row in facts["trust_heads"] if row["expert"] == "seq_v9")
    standalone = scale["seq_v9_standalone_cal_auc"]
    slide_header(
        s,
        "seq_v9 answers when no broker has a mature opinion yet.",
        "A compact sequence expert gives the meta-layer a broad local signal.",
        "All values on this slide are read from facts.scale and facts.trust_heads.",
    )
    add_rect(s, 0.75, 2.05, 3.25, 3.85, fill=NAVY)
    add_textbox(s, 1.0, 2.3, 2.75, 0.38, "Sequence expert", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(
        s,
        1.0,
        3.0,
        2.75,
        1.85,
        "truncated light curve\n+\ntime-aware hidden state\n+\nstandalone calibration",
        size=16,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(s, 1.0, 5.35, 2.75, 0.28, "GRU expert", size=11, bold=True, color=SOFT, align=PP_ALIGN.CENTER)
    add_arrow(s, 4.2, 3.95, 5.2, 3.95, color=GREY, width=2)
    metrics = [
        (5.45, 1.95, whole(scale["seq_v9_params"]), "parameters", "compact GRU"),
        (8.95, 1.95, whole(scale["objects"]), "objects", f"{whole(scale['ztf_objects'])} ZTF / {whole(scale['lsst_objects'])} LSST"),
        (5.45, 3.35, num(standalone["snia_ndet5"]), "standalone cal AUC", "SN Ia @ n_det=5"),
        (8.95, 3.35, num(standalone["other_ndet5"]), "standalone cal AUC", "other @ n_det=5"),
        (5.45, 4.75, whole(seq["n_test"]), "widest test coverage", "seq_v9 trust-head test rows"),
        (8.95, 4.75, whole(scale["stage_a_rows"]), "stage-A rows", "self-supervised sequence scale"),
    ]
    for x, y, value, label, context in metrics:
        add_metric(s, x, y, 3.0, 0.98, value, label, context, color=RED if "coverage" in label else NAVY)
    footer(s)


def add_gates_slide():
    s = prs.slides.add_slide(BLANK)
    slide_header(
        s,
        "The component gates rejected more than they kept.",
        "Credibility comes from reporting what failed, not only what made the final stack.",
        facts["provenance"]["component_gates"],
    )
    rows = []
    for row in facts["component_gates"]:
        delta = row["delta_macro_auc"]
        rows.append(
            (
                row["component"],
                row["decision"],
                "guardrail" if delta is None else num(delta, plus=True),
                ci_text(row["ci95"]),
            )
        )
    add_table(
        s,
        0.65,
        1.9,
        12.0,
        4.75,
        ["Component", "Decision", "Delta macro AUC", "Bootstrap CI / note"],
        rows,
        font_size=12,
        header_size=12,
        col_widths=[2.8, 1.3, 1.65, 3.55],
        num_cols={2},
        first_col_bold=True,
    )
    add_textbox(
        s,
        0.75,
        6.82,
        11.8,
        0.25,
        "Kept components had positive or required calibration value; dropped components did not clear the cal-decided gates.",
        size=10,
        color=GREY,
        italic=True,
        align=PP_ALIGN.CENTER,
    )


def add_deliverables_slide():
    s = prs.slides.add_slide(BLANK)
    slide_header(
        s,
        "Deliverables are ready for science-facing triage, with clear next tests.",
        "The hackathon result is a trust-aware meta-layer plus goal-conditioned ranking outputs.",
        "Next steps focus on stronger truth, out-of-fold experts, and host-aware features.",
    )
    left = [
        (
            "Goal-conditioned priority lists",
            "SN Ia / non-Ia / other users can rank candidates with the same calibrated probability table.",
        ),
        (
            "Conformal abstention",
            "Priority output can say when the model should not make a narrow claim.",
        ),
        (
            "Per-expert confidence payload",
            f"{whole(facts['scale']['experts_registered'])} registered experts; {whole(facts['scale']['experts_trust_headed'])} trust-headed experts.",
        ),
    ]
    right = [
        ("ELAsTiCC2 spectro-truth benchmark", "Benchmark against a cleaner external truth table."),
        ("OOF self-trained experts", "Keep expert training and meta-training separated."),
        ("Host-separation features", "Add host-aware information without weakening temporal safety."),
    ]
    add_textbox(s, 0.7, 1.95, 5.8, 0.36, "Delivered", size=16, bold=True, color=NAVY)
    add_textbox(s, 6.85, 1.95, 5.8, 0.36, "Next", size=16, bold=True, color=NAVY)
    y = 2.48
    for title, body in left:
        bullet_block(s, 0.7, y, 5.8, 1.15, title, body, fill=GREEN_LIGHT, title_color=GREEN)
        y += 1.38
    y = 2.48
    for title, body in right:
        bullet_block(s, 6.85, y, 5.8, 1.15, title, body, fill=LIGHT)
        y += 1.38
    add_rect(s, 0.7, 6.66, 11.95, 0.32, fill=NAVY)
    add_textbox(
        s,
        0.9,
        6.705,
        11.55,
        0.18,
        "Primary product: expert_confidence. Secondary product: trust-weighted follow-up priority.",
        size=8.8,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def build_deck():
    require_assets()
    add_title_slide()
    add_problem_slide()
    add_architecture_slide()
    add_honesty_slide()
    add_headline_slide()
    add_early_epoch_slide()
    add_dp1_slide()
    add_trust_slide()
    add_sequence_expert_slide()
    add_gates_slide()
    add_deliverables_slide()
    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_deck()
