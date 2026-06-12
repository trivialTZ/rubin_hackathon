#!/usr/bin/env python3
"""Build the metaDEBASS v7 status presentation — plain version (mostly tables).

Audience: astronomers and stakeholders unfamiliar with the project.
Run:
    python presentations/build_v7_status.py
Output:
    presentations/metaDEBASS_status_2026-04-26.pptx

Numbers are pulled from the SCC v7 gold table
(data/gold/object_epoch_snapshots_safe_v7.parquet, 211,392 rows / 12,772 objects)
and v7 trust + followup metrics on 2026-04-26.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------- palette (kept minimal) ----------
BLACK = RGBColor(0x00, 0x00, 0x00)
NAVY = RGBColor(0x10, 0x2A, 0x43)
SLATE = RGBColor(0x33, 0x44, 0x55)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
SOFT = RGBColor(0xCB, 0xD5, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x1D, 0x4E, 0xD8)
GREEN = RGBColor(0x05, 0x80, 0x4F)


# ---------- presentation ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def add_textbox(slide, x, y, w, h, text, *,
                size=14, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = "Helvetica"
    return tb


def add_rect(slide, x, y, w, h, fill=LIGHT, line=None, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def slide_header(slide, title, subtitle=None):
    add_textbox(slide, 0.5, 0.35, 12.3, 0.7, title,
                size=26, bold=True, color=BLACK)
    if subtitle:
        add_textbox(slide, 0.5, 1.0, 12.3, 0.5, subtitle,
                    size=14, color=GREY, italic=True)
    # thin underline
    add_rect(slide, 0.5, 1.45, 12.3, 0.02, fill=SLATE)


def add_table(slide, x, y, w, h, headers, rows, *,
              header_fill=NAVY, header_color=WHITE,
              font_size=12, header_size=12,
              col_widths=None, num_cols=None,
              first_col_bold=False):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x),
                                       Inches(y), Inches(w), Inches(h))
    tbl = tbl_shape.table

    if col_widths is not None:
        total = sum(col_widths)
        for ci, frac in enumerate(col_widths):
            tbl.columns[ci].width = Inches(w * frac / total)

    num_cols = num_cols or set()

    for ci, htext in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        p.text = htext
        p.font.size = Pt(header_size)
        p.font.bold = True
        p.font.color.rgb = header_color
        p.font.name = "Helvetica"
        p.alignment = PP_ALIGN.RIGHT if ci in num_cols else PP_ALIGN.LEFT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.text = str(val)
            p.font.size = Pt(font_size)
            p.font.color.rgb = BLACK
            p.font.name = "Helvetica"
            p.font.bold = first_col_bold and ci == 0
            p.alignment = PP_ALIGN.RIGHT if ci in num_cols else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
    return tbl


def add_arrow(slide, x1, y1, x2, y2, color=BLACK, width=1.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    line = conn.line
    lineFmt = line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree
    tail = etree.SubElement(lineFmt, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return conn


# ============================================================
# Slide 1 — Title (very plain)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_textbox(s, 0.5, 2.4, 12.3, 1.2, "metaDEBASS",
            size=72, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 3.6, 12.3, 0.8,
            "Trust-aware expert confidence for early Rubin/LSST transient epochs",
            size=22, color=SLATE, align=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 4.7, 12.3, 0.4,
            "Status: v7  ·  April 26, 2026",
            size=16, color=GREY, align=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 5.3, 12.3, 0.4,
            "Xianzhe (TZ) Tang  ·  Boston University",
            size=14, color=GREY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 2 — Outline  (Winston: open with a promise + map)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "What you'll learn in this talk",
             subtitle="A 6-section map.  The right column points to each section.")

headers = ["#", "Section", "Pages"]
rows = [
    ("1", "The problem  —  pick the right ~100 transients to spectrum from ~10⁴ candidates per LSST night",  "3"),
    ("2", "The approach  —  preserve timing, calibrate trust, then score follow-up as a downstream proxy",    "4"),
    ("3", "The model  —  per-expert confidence + optional follow-up proxy + monotonic remapping",             "5–6"),
    ("4", "Training data + truth  —  dense ZTF transferred to sparse LSST DP1",                               "7"),
    ("5", "Where metaDEBASS sits  —  vs single-broker, photometric typing, BTSbot",                          "8"),
    ("6", "Status  —  12 of 28 calibrated today; secondary follow-up proxy AUC 0.970",                         "9–20"),
]
add_table(s, 0.5, 1.7, 12.3, 4.4, headers, rows,
          font_size=14, header_size=14,
          col_widths=[0.4, 10.4, 1.5], num_cols={0, 2})

add_textbox(s, 0.5, 6.3, 12.3, 0.4,
            "The promise",
            size=14, bold=True, color=BLACK)
add_textbox(s, 0.5, 6.75, 12.3, 0.7,
            "By the end you should know how metaDEBASS takes raw outputs from 7 brokers and turns them into "
            "calibrated per-expert confidence at each object epoch, plus a secondary follow-up proxy for triage.",
            size=12, color=SLATE)


# ============================================================
# Slide 3 (intro) — The problem
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "The problem we are solving",
             subtitle="LSST will publish more transients per night than any team can spectroscopically follow up — "
                      "we have to pick the right ones early.")

left_x, right_x = 0.5, 6.95
box_w, box_h, top_y = 5.85, 5.4, 1.7

add_rect(s, left_x, top_y, box_w, box_h, fill=LIGHT, line=SOFT)
add_textbox(s, left_x + 0.2, top_y + 0.15, box_w - 0.4, 0.45,
            "The alert funnel", size=18, bold=True, color=NAVY)
add_textbox(s, left_x + 0.2, top_y + 0.7, box_w - 0.4, 4.6,
            "•  ~10⁷  alerts per night, Rubin southern sky\n\n"
            "•  ~10⁵  real transients\n"
            "    (after asteroid / variable-star / instrumental rejection)\n\n"
            "•  ~10⁴  SN-like candidates worth study\n\n"
            "•  ~10²  candidates can actually receive spectra\n"
            "    (combined global capacity, per night)\n\n"
            "→  We have to pick the right ~100 from ~10,000 — every night.",
            size=14, color=BLACK)

add_rect(s, right_x, top_y, box_w, box_h, fill=LIGHT, line=SOFT)
add_textbox(s, right_x + 0.2, top_y + 0.15, box_w - 0.4, 0.45,
            "The constraint", size=18, bold=True, color=NAVY)
add_textbox(s, right_x + 0.2, top_y + 0.7, box_w - 0.4, 4.6,
            "•  Decision time = after 3–5 detections,\n"
            "    before the light-curve is mature\n\n"
            "•  Each broker (ALeRCE, Fink, Lasair, ANTARES,\n"
            "    Pitt-Google, AMPEL, Babamul) publishes its own\n"
            "    classifications — different strengths, different\n"
            "    blind spots, different output spaces\n\n"
            "•  We need calibrated per-expert confidence:\n"
            "    which expert is available, trustworthy, and\n"
            "    contributing useful evidence at this epoch?\n\n"
            "→  metaDEBASS emits expert_confidence first;\n"
            "    p_follow_proxy is a secondary triage score.",
            size=14, color=BLACK)


# ============================================================
# Slide 4 (intro) — metaDEBASS approach (3-stage pipeline)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "The metaDEBASS approach — ingest, calibrate, report trust",
             subtitle="Three stages landing in expert_confidence first, with a downstream follow-up proxy.")

stage_w = 3.9
stage_h = 3.6
stage_y = 1.8
stage_xs = [0.5, 4.7, 8.9]
labels = ["1.  INGEST", "2.  CALIBRATE", "3.  SCORE"]
contents = [
    "•  7 broker services  —  pulled via REST API\n"
    "    or Google BigQuery\n"
    "•  6 classifiers we re-run locally on the\n"
    "    cached light-curve  (SuperNNova, SALT3,\n"
    "    ALeRCE-LC, ParSNIP, Bazin/Villar, ORACLE)\n\n"
    "Raw-event table  (one row per broker event,\n"
    "deduplicated)\n"
    "        ↓\n"
    "Object × epoch table  (one row per object\n"
    "per detection step;  features built from\n"
    "ONLY the first N detections — never future)",

    "•  12 small reliability models —\n"
    "    one per classifier.\n\n"
    "•  Each one learns from data:\n"
    "    given the light-curve so far + which\n"
    "    survey, how often does THIS classifier\n"
    "    get it right?\n\n"
    "•  Two flavours, matching what the\n"
    "    classifier outputs:\n"
    "    –  SN-vs-not screens   (e.g. SLSN-RF)\n"
    "    –  Type-Ia callers     (most others)",

    "•  Primary output:\n"
    "    expert_confidence per expert.\n\n"
    "•  Inputs:\n"
    "    –  the 12 calibrated trust scores\n"
    "    –  51 light-curve numbers  (u,g,r,i,z,y)\n"
    "    –  Babamul context  (known star?  rock?\n"
    "        cross-match to another survey?)\n"
    "    –  Lasair host-galaxy info\n\n"
    "•  Secondary output:\n"
    "    p_follow_proxy — is the candidate\n"
    "    worth a spectrum?",
]
for x, label, content in zip(stage_xs, labels, contents):
    add_rect(s, x, stage_y, stage_w, stage_h, fill=LIGHT, line=SOFT)
    add_rect(s, x, stage_y, stage_w, 0.55, fill=NAVY)
    add_textbox(s, x, stage_y + 0.07, stage_w, 0.4,
                label, size=16, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.15, stage_y + 0.7, stage_w - 0.3,
                stage_h - 0.85, content, size=11, color=BLACK)

arrow_y = stage_y + stage_h / 2
add_arrow(s, stage_xs[0] + stage_w + 0.05, arrow_y,
          stage_xs[1] - 0.05, arrow_y, color=GREY, width=2.0)
add_arrow(s, stage_xs[1] + stage_w + 0.05, arrow_y,
          stage_xs[2] - 0.05, arrow_y, color=GREY, width=2.0)

add_rect(s, 0.5, 5.85, 12.3, 1.0, fill=NAVY)
add_textbox(s, 0.5, 5.95, 12.3, 0.4,
            "v7 result", size=14, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 6.35, 12.3, 0.5,
            "211,392 detection-step rows × 12,772 objects   ·   "
            "follow-up proxy AUC 0.970   ·   calibrated ECE 0.006",
            size=15, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 5 (intro) — Model architecture
# ============================================================
s = prs.slides.add_slide(BLANK)

fig_path = Path(__file__).resolve().parent / "figures" / "architecture_v7.png"
if fig_path.exists():
    s.shapes.add_picture(str(fig_path),
                         Inches(0.42), Inches(0.15),
                         height=Inches(7.0))
else:
    add_rect(s, 0.45, 1.55, 12.3, 4.9, fill=LIGHT, line=SOFT)
    add_textbox(s, 0.45, 3.7, 12.3, 0.5,
                "[architecture_v7.png — generate via figures/_gen_architecture.py]",
                size=12, color=GREY, align=PP_ALIGN.CENTER)

add_textbox(s, 0.5, 7.18, 12.3, 0.16,
            "Machine-readable contract: expert_confidence primary; p_follow_proxy secondary; latest_object_unsafe excluded; 51 LC features; 28 experts; 12 calibrated v7 heads.",
            size=6, color=GREY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 6 (intro) — Calibration before/after  (uses reliability_v7.png)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Follow-up proxy calibration:  ECE  0.033  →  0.006",
             subtitle="Isotonic regression sharpens downstream probabilities while preserving ranking.")

fig_path = Path(__file__).resolve().parent / "figures" / "reliability_v7.png"
if fig_path.exists():
    s.shapes.add_picture(str(fig_path),
                         Inches(0.5), Inches(1.7),
                         width=Inches(7.0))
else:
    add_rect(s, 0.5, 1.7, 7.0, 5.0, fill=LIGHT, line=SOFT)
    add_textbox(s, 0.5, 4.0, 7.0, 0.5,
                "[reliability_v7.png — generate via figures/_gen_reliability.py]",
                size=12, color=GREY, align=PP_ALIGN.CENTER)

add_textbox(s, 7.7, 1.75, 5.2, 0.4,
            "Why it matters", size=16, bold=True, color=NAVY)
add_textbox(s, 7.7, 2.2, 5.2, 4.9,
            "The primary metaDEBASS payload is still\n"
            "expert_confidence: which experts are\n"
            "available, trustworthy, and informative\n"
            "at this object epoch.\n\n"
            "The follow-up proxy is the downstream\n"
            "operational score for triage: P(worth a\n"
            "spectrum).\n\n"
            "A well-calibrated 0.9 should mean: in 9 out\n"
            "of 10 such cases, it really is worth one.\n\n"
            "On the v7 held-out test split:\n\n"
            "    ECE 0.033 raw  →  0.006 calibrated\n"
            "    AUC 0.9701 raw →  0.9700 calibrated\n\n"
            "The ranking is essentially unchanged;\n"
            "calibration improves the meaning of the\n"
            "probability scale.",
            size=12, color=BLACK)


# ============================================================
# Slide 7 (intro) — Training data + truth
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Training data and ground truth",
             subtitle="Mixed ZTF + LSST cohort. Five-tier truth for label coverage.")

add_textbox(s, 0.5, 1.7, 6.0, 0.4,
            "Cohort (v7 gold)", size=16, bold=True, color=NAVY)
headers = ["Survey", "Objects", "Epoch rows"]
rows = [
    ("ZTF",        "8,774",   "139,061"),
    ("LSST",       "3,998",   "72,331"),
    ("Combined",   "12,772",  "211,392"),
]
add_table(s, 0.5, 2.15, 6.0, 2.6, headers, rows,
          font_size=12, header_size=12,
          col_widths=[1.6, 1.6, 2.8], num_cols={1, 2})

add_textbox(s, 6.7, 1.7, 6.0, 0.4,
            "Truth — 5 tiers  (best to weakest)", size=16, bold=True, color=NAVY)
headers = ["Tier", "Source", "Quality"]
rows = [
    ("1", "TNS spectroscopic type",                                "strong  (spec.)"),
    ("2", "TNS name only  (no spec type yet)",                     "name only"),
    ("3", "Broker consensus  (≥ 2 classifiers agree)",             "consensus"),
    ("4", "Host context  (SIMBAD galaxy + low star probability)",  "context"),
    ("5", "ALeRCE discovery stamp class",                          "weak"),
]
add_table(s, 6.7, 2.15, 6.1, 2.9, headers, rows,
          font_size=12, header_size=12,
          col_widths=[0.5, 4.4, 1.2])

add_textbox(s, 0.5, 5.4, 12.3, 0.4,
            "Why we need labels from many sources",
            size=14, bold=True, color=BLACK)
add_textbox(s, 0.5, 5.85, 12.3, 1.5,
            "Rubin DP1 (the first preview, Feb 2026) has only ~15 spectroscopically-confirmed SNe across the published papers "
            "(Dong+, Aleo+, Smith+) — too few to train per-classifier reliability models on LSST data alone.\n\n"
            "Solution:  train on the ZTF arm, where the Transient Name Server (TNS) gives us thousands of spec-confirmed labels, "
            "then transfer the calibrated model to LSST. The 'which survey' flag lets the model adapt to LSST automatically — "
            "no hand-coded rules.",
            size=12, color=SLATE)


# ============================================================
# Slide 8 (intro) — Where metaDEBASS sits in the broker landscape
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Where metaDEBASS sits in the broker landscape",
             subtitle="Not another physics classifier — a meta-classifier that combines what brokers already produce.")

headers = ["Approach", "Where it fits", "Limitation"]
rows = [
    ("Single broker  (e.g. Fink SNN, ALeRCE LC)",
     "Specialist per-broker classification",
     "Misses signal from other brokers; no cross-broker calibration"),
    ("Photometric typing model  (SuperNNova, ParSNIP, RAPID)",
     "Per-class probabilities from light-curve fit",
     "Needs mature LC; doesn't combine multi-source evidence"),
    ("Targeted SN follow-up trigger  (BTSbot, Lasair-DH)",
     "Operational follow-up tool, hand-tuned per survey",
     "Survey-specific; not a general meta-classifier"),
    ("Babamul  (broker-of-brokers)",
     "Live event re-broadcasting + cross-broker tags (star/rock/x-match)",
     "Tags are categorical context; doesn't combine classifier probabilities"),
    ("metaDEBASS  (this work)",
     "Reports calibrated expert_confidence for each object epoch, with a secondary follow-up proxy after 3-5 detections",
     "Calibration depends on the training cohort; tied to broker uptime"),
]
add_table(s, 0.5, 1.7, 12.3, 3.6, headers, rows,
          font_size=12, header_size=12,
          col_widths=[3.0, 4.5, 4.8])

add_textbox(s, 0.5, 5.55, 12.3, 0.4,
            "What's distinctive about metaDEBASS",
            size=14, bold=True, color=BLACK)
add_textbox(s, 0.5, 5.95, 12.3, 1.5,
            "•  Per-expert trust is calibrated — each available classifier has an explicit confidence value\n"
            "•  Survey-agnostic — the same pipeline trains on ZTF and scores on LSST\n"
            "•  Causal-safe by construction — when scoring at detection N, the model sees ONLY the first N detections\n"
            "•  Built for very early decisions — expert_confidence is emitted before a mature light-curve",
            size=12, color=BLACK)


# ============================================================
# Slide 9 — Headline counts
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "v7 — 12 calibrated trust heads across 7 brokers",
             subtitle="All numbers measured on the v7 dataset — "
                      "211,392 detection-step rows / 12,772 unique objects.")

headers = ["Quantity", "Count"]
rows = [
    ("Brokers connected",                          "7"),
    ("Classifiers registered (have an adapter)",   "28"),
    ("Classifiers producing data on v7 cohort",    "15"),
    ("Classifiers calibrated and used in v7",      "12"),
    ("Light-curve features (per epoch, ugrizy)",   "51"),
    ("Secondary follow-up proxy ranking AUC",                   "0.970"),
    ("Secondary follow-up proxy ECE",                           "0.006"),
    ("Primary scored payload",                                  "expert_confidence[expert_key]"),
    ("Follow-up calibration step (raw → calibrated ECE)",       "0.033 → 0.006"),
]
add_table(s, 0.5, 1.7, 12.3, 5.0, headers, rows,
          font_size=16, header_size=15,
          col_widths=[3.0, 1.5], num_cols={1},
          first_col_bold=False)

add_textbox(s, 0.5, 6.85, 12.3, 0.4,
            "Sources: data/gold/object_epoch_snapshots_safe_v7.parquet  +  "
            "reports/metrics/{expert_trust,followup}_metrics_safe_v7.json",
            size=10, color=GREY, italic=True)


# ============================================================
# Slide 10 — The 7 brokers (one row per broker)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "7 brokers + 6 local re-runners feed one common gold table",
             subtitle="Each broker exposes one or more classifiers.")

headers = ["Broker", "Origin", "Access mechanism",
           "Classifiers registered", "Calibrated in v7"]
rows = [
    ("ALeRCE",      "Chile",         "REST API",                "7", "0  (snapshot mode — see slide 5)"),
    ("Fink",        "France",        "REST API",                "6", "6"),
    ("Lasair",      "UK",            "REST API",                "1", "0  (used as feature)"),
    ("Pitt-Google", "USA",           "Google BigQuery",         "3", "1"),
    ("ANTARES",     "NOIRLab (USA)", "REST API",                "2", "0"),
    ("AMPEL",       "Germany / Switzerland", "Vendored — runs locally", "2", "1"),
    ("Babamul",     "Broker-of-brokers", "REST API  (Kafka verified, not wired)", "1", "0  (used as feature)"),
]
add_table(s, 0.5, 1.7, 12.3, 3.5, headers, rows,
          font_size=14, header_size=14,
          col_widths=[1.4, 1.7, 2.4, 1.5, 1.7],
          num_cols={3, 4})

# Local rerunners block
add_textbox(s, 0.5, 5.4, 12.3, 0.4,
            "Plus: 6 local re-runnable classifiers — we execute the same upstream models on our cached light curves.",
            size=14, color=BLACK)
add_table(s, 0.5, 5.85, 12.3, 1.0,
          ["Source", "What it is", "Classifiers", "Calibrated in v7"],
          [("Local re-run", "SuperNNova / ParSNIP / ALeRCE-LC / SALT3 χ² / Bazin-Villar / ORACLE",
            "6", "4")],
          font_size=14, header_size=14,
          col_widths=[1.4, 5.6, 1.5, 1.7], num_cols={2, 3})


# ============================================================
# Slide 11 — How information arrives (small diagram + bullets)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Three ingest paths converge in silver, then in gold",
             subtitle="REST, BigQuery, and local execution all land in one common table.")

# Diagram: 3 input boxes -> silver -> gold
inputs = [
    ("REST polling",     "ALeRCE · Fink · Lasair\nANTARES · Babamul"),
    ("Google BigQuery",  "Pitt-Google"),
    ("Local execution",  "AMPEL SNGuess\n6 local re-runners"),
]

box_w, box_h = 2.6, 1.35
gap = 0.3
top = 1.85
for i, (h, body) in enumerate(inputs):
    y = top + i * (box_h + gap)
    add_rect(s, 0.7, y, box_w, box_h, fill=LIGHT, line=SOFT, rounded=False)
    add_textbox(s, 0.85, y + 0.12, box_w - 0.3, 0.4, h,
                size=14, bold=True, color=BLACK)
    add_textbox(s, 0.85, y + 0.5, box_w - 0.3, 0.8, body,
                size=11, color=SLATE)

# arrows to silver
mid_x = 4.4
silver_x = 5.0
silver_y = 2.85
silver_w = 3.4
silver_h = 1.4
add_rect(s, silver_x, silver_y, silver_w, silver_h, fill=NAVY)
add_textbox(s, silver_x, silver_y + 0.2, silver_w, 0.5,
            "Raw-event table", size=18, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, silver_x + 0.1, silver_y + 0.7, silver_w - 0.2, 0.6,
            "one row per broker event\n(deduplicated, normalised)",
            size=11, color=WHITE, align=PP_ALIGN.CENTER)

for i in range(3):
    y = top + i * (box_h + gap) + box_h / 2
    add_arrow(s, 0.7 + box_w + 0.05, y, silver_x - 0.05,
              silver_y + silver_h / 2, color=GREY, width=1.0)

# arrow to gold
gold_x = 9.4
gold_y = silver_y
gold_w = 3.4
gold_h = silver_h
add_rect(s, gold_x, gold_y, gold_w, gold_h, fill=NAVY)
add_textbox(s, gold_x, gold_y + 0.2, gold_w, 0.5,
            "Object × epoch table", size=18, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, gold_x + 0.1, gold_y + 0.7, gold_w - 0.2, 0.7,
            "one row per object × per detection step\n51 light-curve features + 28 classifier columns\n(features built from first N detections only)",
            size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_arrow(s, silver_x + silver_w + 0.05, silver_y + silver_h / 2,
          gold_x - 0.05, gold_y + gold_h / 2, color=GREY, width=1.5)

# bullets below
add_textbox(s, 0.5, 6.05, 12.3, 0.4,
            "Notes",
            size=14, bold=True, color=BLACK)
add_textbox(s, 0.5, 6.45, 12.3, 1.0,
            "•  Each event stores the broker's score AND its publish timestamp, so we can match "
            "scores to detections without ever pulling in a future score.\n"
            "•  Local re-runners (SuperNNova, SALT3, etc.) execute directly on the cached light-curve — "
            "no broker round-trip — so they are available for almost every object.",
            size=12, color=SLATE)


# ============================================================
# Slide 12 — ALeRCE classifiers
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "ALeRCE — REST works, but the snapshot scores can't be used at training time",
             subtitle="Chilean broker. 351,594 ALeRCE events arriving across 14 upstream keys; 143,643 are tagged with our "
                      "7 registered classifiers. ALL are blocked from training — explained below.")

headers = ["Classifier key", "Type", "Survey",
           "Silver events", "In v7 gold?", "Calibrated in v7"]
rows = [
    ("alerce/LC_classifier_ATAT_forced_phot(beta)",      "Transformer (beta)",        "ZTF",  "68,471", "blocked (training safety)", "no"),
    ("alerce/stamp_classifier",                          "Image stamp (legacy)",      "ZTF",  "36,105", "blocked (training safety)", "no"),
    ("alerce/stamp_classifier_2025_beta",                "Image stamp (2025 beta)",   "ZTF",  "11,766", "blocked (training safety)", "no"),
    ("alerce/lc_classifier_BHRF_forced_phot_transient",  "Random forest (BHRF)",      "ZTF",  "10,974", "blocked (training safety)", "no"),
    ("alerce/lc_classifier_transient",                   "Light-curve ML (legacy)",   "ZTF",  "10,840", "blocked (training safety)", "no"),
    ("alerce/lc_classifier_BHRF_forced_phot_top",        "BHRF top-level",            "ZTF",  "5,487",  "blocked (training safety)", "no"),
    ("alerce/stamp_classifier_rubin_beta",               "Image stamp (Rubin beta)",  "LSST", "0",      "no events yet",         "no"),
]
add_table(s, 0.5, 1.7, 12.3, 3.7, headers, rows,
          font_size=12, header_size=12,
          col_widths=[3.6, 2.0, 0.9, 1.4, 2.4, 1.4], num_cols={3, 4, 5})

add_textbox(s, 0.5, 5.65, 12.3, 0.4,
            "Why every row says 'blocked'",
            size=14, bold=True, color=BLACK)
add_textbox(s, 0.5, 6.05, 12.3, 1.4,
            "•  ALeRCE's REST API returns ONE current score per object — not a separate score per detection. "
            "We don't know WHEN that score was computed.\n"
            "•  If we trained on it, the model would learn from classifications produced AFTER the cutoff we're scoring at — "
            "i.e. the future would leak into the past, inflating accuracy artificially.\n"
            "•  Workaround in v7:  re-run the same ALeRCE LC model locally on the cached light-curve, properly truncated to "
            "the first N detections. That covers 100% of the cohort with AUC 0.88  (slide 16).\n"
            "•  At scoring time (not training time) the snapshot scores ARE useful — that's a future toggle.",
            size=12, color=SLATE)


# ============================================================
# Slide 13 — Fink classifiers
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Fink — 6 of 6 calibrated; the broadest broker contribution",
             subtitle="French broker. Cleanest per-alert API; carries the most calibrated classifiers in v7.")

headers = ["Classifier key", "Type", "Survey",
           "Coverage on v7 cohort", "AUC", "Calibrated in v7"]
rows = [
    ("fink/snn",                "SuperNNova SN-Ia probability",      "ZTF",  "56.1%", "0.92", "yes"),
    ("fink/rf_ia",              "Random-forest SN-Ia probability",   "ZTF",  "56.1%", "0.81", "yes"),
    ("fink/slsn   (NEW v7)",    "SLSN-RF (super-luminous SN screen)","ZTF",  "9.3%",  "0.93", "yes"),
    ("fink_lsst/snn",           "SuperNNova SN-Ia probability",      "LSST", "27.4%", "0.84", "yes"),
    ("fink_lsst/cats",          "CATS (CNN multi-class)",            "LSST", "23.5%", "0.84", "yes"),
    ("fink_lsst/early_snia",    "Pre-peak SN-Ia trigger",            "LSST", "0.4%",  "0.86", "yes"),
]
add_table(s, 0.5, 1.7, 12.3, 4.0, headers, rows,
          font_size=12, header_size=12,
          col_widths=[2.8, 3.0, 0.9, 1.7, 1.0, 1.4], num_cols={3, 4, 5})

add_textbox(s, 0.5, 6.0, 12.3, 0.4,
            "Notes",
            size=14, bold=True, color=BLACK)
add_textbox(s, 0.5, 6.4, 12.3, 1.0,
            "• Fink LSST EarlySNIa coverage is 0.4% because it requires ≥7 epochs to fire — "
            "this rises sharply as LSST coverage densifies.\n"
            "• fink/slsn coverage is 9.3% because most objects are not SLSN candidates — its calibrated AUC 0.93 "
            "still makes it a strong screening signal where it does fire.",
            size=12, color=SLATE)


# ============================================================
# Slide 14 — Pitt-Google + ANTARES + AMPEL
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "AMPEL SNGuess (NEW v7) is the top-ranked classifier  (AUC 0.95)",
             subtitle="Three more brokers — Pitt-Google delivers via BigQuery; AMPEL is run locally.")

headers = ["Broker", "Classifier key", "Type", "Survey",
           "Coverage on v7 cohort", "AUC", "Calibrated in v7"]
rows = [
    ("Pitt-Google", "supernnova_lsst",       "SuperNNova SN-Ia probability",       "LSST", "18.6%", "0.91", "yes"),
    ("Pitt-Google", "supernnova_ztf",        "SuperNNova SN-Ia probability",       "ZTF",  "0%",    "—",    "no  (7K silver events; BQ ZTF table has no per-alert timestamp → filtered)"),
    ("Pitt-Google", "upsilon_lsst",          "UPSILoN (variable star typing)",     "LSST", "0%",    "—",    "no  (not yet backfilled)"),
    ("ANTARES",     "oracle",                "ORACLE multi-class",                 "both", "0%",    "—",    "no  (upstream not exposed)"),
    ("ANTARES",     "superphot_plus",        "Superphot+ Bayesian fitter",         "both", "0.7%",  "—",    "no  (too sparse to calibrate)"),
    ("AMPEL",       "snguess  (NEW v7)",     "SNGuess (XGBoost; bright SN prob)",  "both", "79.9%", "0.95", "yes  (top-ranked input)"),
    ("AMPEL",       "parsnip_followme",      "ParSNIP follow-me wrapper",          "both", "0%",    "—",    "no  (dormant upstream)"),
]
add_table(s, 0.5, 1.7, 12.3, 4.5, headers, rows,
          font_size=11, header_size=12,
          col_widths=[1.2, 2.2, 2.6, 0.9, 1.5, 0.9, 3.0], num_cols={4, 5})


# ============================================================
# Slide 15 — Lasair + Babamul + Local re-runners
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "4 local re-runners cover ≥ 87% of the cohort",
             subtitle="Two context feeds (Lasair, Babamul) + 6 classifiers we run on our cached light curves.")

headers = ["Source", "Classifier key", "Type", "Survey",
           "Coverage on v7 cohort", "AUC", "Calibrated in v7"]
rows = [
    ("Lasair",      "sherlock",     "Host-galaxy context",                  "both", "34.2%",  "—",   "no  (used as feature on followup head)"),
    ("Babamul",     "babamul",      "Star / asteroid / x-match flags",      "both", "45.8%",  "—",   "no  (used as feature on followup head)"),
    ("Local rerun", "supernnova",   "SuperNNova local",                     "both", "87.0%",  "0.91", "yes"),
    ("Local rerun", "alerce_lc",    "ALeRCE LC classifier (local rerun)",   "both", "100.0%", "0.88", "yes"),
    ("Local rerun", "salt3_chi2",   "SALT3 χ² Ia template-fit consistency", "both", "92.0%",  "0.90", "yes"),
    ("Local rerun", "lc_features_bv","Bazin / Villar parametric features",  "both", "99.4%",  "0.88", "yes"),
    ("Local rerun", "parsnip",      "ParSNIP local",                        "both", "0%",     "—",   "no  (not yet inferred on cohort)"),
    ("Local rerun", "oracle_lsst",  "ORACLE local",                         "LSST", "0%",     "—",   "no  (domain mismatch on DP1)"),
]
add_table(s, 0.5, 1.7, 12.3, 5.0, headers, rows,
          font_size=11, header_size=12,
          col_widths=[1.3, 2.0, 3.0, 0.8, 1.5, 0.9, 2.8], num_cols={4, 5})


# ============================================================
# Slide 16 — The 12 calibrated classifiers (bar chart + key takeaways)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "The 12 calibrated classifiers, ranked AUC  0.95 → 0.81",
             subtitle="Each one has a small per-classifier reliability model trained on the v7 cohort.")

fig_path = Path(__file__).resolve().parent / "figures" / "top12_auc.png"
if fig_path.exists():
    s.shapes.add_picture(str(fig_path),
                         Inches(0.3), Inches(1.65),
                         width=Inches(8.0))
else:
    add_rect(s, 0.3, 1.65, 8.0, 5.0, fill=LIGHT, line=SOFT)
    add_textbox(s, 0.3, 4.0, 8.0, 0.5,
                "[top12_auc.png — generate via figures/_gen_top12_bar.py]",
                size=12, color=GREY, align=PP_ALIGN.CENTER)

add_textbox(s, 8.5, 1.65, 4.5, 0.45,
            "Key takeaways", size=15, bold=True, color=NAVY)
add_textbox(s, 8.5, 2.15, 4.5, 5.0,
            "•  AMPEL SNGuess (NEW v7) tops at\n"
            "    AUC 0.95 — highest of any single\n"
            "    classifier in the deck.\n\n"
            "•  Fink contributes 6 of the top 12 —\n"
            "    broadest single-broker contribution.\n\n"
            "•  Local re-runners (4 of 12) bridge\n"
            "    the LSST coverage gap while LSST\n"
            "    broker outputs mature.\n\n"
            "•  fink_lsst/early_snia  (AUC 0.86)\n"
            "    is small-coverage but high signal —\n"
            "    fires only when ≥ 7 epochs available.\n\n"
            "•  Their calibrated q_i values populate\n"
            "    expert_confidence; the downstream\n"
            "    follow-up proxy reaches AUC 0.970.",
            size=12, color=BLACK)

add_textbox(s, 0.5, 6.85, 12.3, 0.4,
            "The trust heads are the primary science product; the follow-up proxy is the secondary triage layer.",
            size=13, color=BLACK, italic=True, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 17 — Coverage by survey (counts table)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "LSST  4/7 calibrated  ·  ZTF  3/10  ·  Survey-agnostic  5/11",
             subtitle="ZTF is the mature stream; LSST is at DP1 — first preview only.")

headers = ["Survey", "Classifiers registered", "Producing data on v7 cohort",
           "Calibrated in v7"]
rows = [
    ("ZTF only",                 "10",  "3",  "3"),
    ("LSST only",                 "7",  "4",  "4"),
    ("Survey-agnostic + local",  "11",  "8",  "5"),
    ("TOTAL",                    "28", "15", "12"),
]
add_table(s, 0.5, 1.7, 12.3, 3.5, headers, rows,
          font_size=15, header_size=14,
          col_widths=[2.4, 2.4, 3.0, 2.0], num_cols={1, 2, 3})

# Notes table
add_textbox(s, 0.5, 5.4, 12.3, 0.4,
            "What the surveys give us today",
            size=14, bold=True, color=BLACK)
headers = ["Survey", "What works today", "Main limitation"]
rows = [
    ("ZTF",  "Fink ×3, AMPEL SNGuess, 4 local re-runners — dense coverage",
             "ALeRCE REST is snapshot-only (no per-alert scores); covered by local rerun"),
    ("LSST", "Fink ×3, PGB SuperNNova, AMPEL SNGuess, 4 local re-runners",
             "Coverage is sparse on DP1; some classifiers need ≥7 epochs to fire"),
]
add_table(s, 0.5, 5.85, 12.3, 1.4, headers, rows,
          font_size=11, header_size=12,
          col_widths=[1.0, 5.6, 5.7])


# ============================================================
# Slide 18 — Why the other 16 are not calibrated yet
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Why are 16 of 28 classifiers not yet calibrated?",
             subtitle="Six distinct reasons — none of them are blocking the v7 product.  "
                      "Sum of 'Count' column = 16.")

headers = ["Reason", "Count", "Examples", "Plan"]
rows = [
    ("Used as a context feature, not a stand-alone classifier",
     "2",
     "Lasair Sherlock, Babamul",
     "Stays as a feature on the follow-up proxy — never gets its own trust head"),
    ("Snapshot-only output  →  blocked at training time (causality)",
     "7",
     "6 ALeRCE classifiers (143K events) + Pitt-Google SuperNNova ZTF (7K events)",
     "Local ALeRCE re-run covers ALeRCE; PGB ZTF needs an upstream timestamp"),
    ("No events arriving yet for our cohort",
     "2",
     "ALeRCE Rubin stamp (LSST), PGB UPSILoN LSST",
     "Backfill in v8; both are upstream-ready"),
    ("Upstream classifier not yet exposed by the broker",
     "3",
     "ANTARES ORACLE, AMPEL ParSNIP-FollowMe, local ParSNIP",
     "Wait for the broker to publish; turn on the adapter when available"),
    ("Trained on simulated LSST  →  doesn't transfer to real DP1",
     "1",
     "Local ORACLE-LSST",
     "Retrain on real LSST labels (available after Rubin DR1)"),
    ("Too few examples to fit a reliability model  (< 1% coverage)",
     "1",
     "ANTARES Superphot+ (0.7% coverage)",
     "Re-evaluate after broader backfill"),
]
add_table(s, 0.5, 1.7, 12.3, 5.0, headers, rows,
          font_size=12, header_size=12,
          col_widths=[3.4, 0.7, 3.6, 4.6], num_cols={1})


# ============================================================
# Slide 19 — Roadmap (counts over time)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Calibrated count: 12 today → ~18 by Rubin DR1 (~2028)",
             subtitle="Realistic projection assuming brokers ship what they have already announced.")

headers = ["Date",            "Calibrated count", "What gets added",                                              "Trigger"]
rows = [
    ("Today (Apr 2026)",        "12",  "(baseline v7)",                                              "—"),
    ("Late 2026 — LSST commissioning ends",
                                "14",  "Fink LSST EarlySNIa firing routinely,  ALeRCE Rubin stamp backfill",
                                                                                                     "More LSST coverage per object"),
    ("Mid 2027 — full survey + brokers catch up",
                                "16",  "ALeRCE LSST-ATAT release,  Babamul live LSST stream matures",
                                                                                                     "Upstream releases"),
    ("DR1 (~2028)  — first full Rubin Data Release",
                                "≈ 18", "ANTARES LSST classifiers,  ORACLE retrain on real LSST labels,  PGB ZTF timestamp fix",
                                                                                                     "Real LSST labels available for training"),
]
add_table(s, 0.5, 1.7, 12.3, 4.5, headers, rows,
          font_size=12, header_size=12,
          col_widths=[3.0, 1.5, 5.5, 2.3], num_cols={1})

add_textbox(s, 0.5, 6.3, 12.3, 0.4,
            "Caveat",
            size=14, bold=True, color=BLACK)
add_textbox(s, 0.5, 6.7, 12.3, 0.6,
            "Counts above are conservative — they only count classifiers we can already ingest. "
            "New broker classifiers announced during 2026–2028 (e.g. Babamul ACAI, BTSbot, Fink TDE) would be additional.",
            size=12, color=SLATE, italic=True)


# ============================================================
# Slide 20 — Summary
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "metaDEBASS v7 — calibrated expert confidence at early epochs")

headers = ["", ""]
rows = [
    ("Brokers connected",                     "7"),
    ("Classifiers registered",                "28"),
    ("Classifiers producing data in v7",      "15"),
    ("Classifiers calibrated and used in v7", "12"),
    ("Primary output",                        "expert_confidence"),
    ("Secondary follow-up proxy AUC",         "0.970"),
    ("Secondary follow-up proxy ECE",         "0.006"),
    ("Realistic count by DR1 (~2028)",        "≈ 18"),
]
add_table(s, 0.5, 1.7, 12.3, 4.0, headers, rows,
          font_size=18, header_size=12,
          col_widths=[3.5, 1.0], num_cols={1})

add_textbox(s, 0.5, 6.0, 12.3, 0.4,
            "One-liner",
            size=14, bold=True, color=BLACK)
add_textbox(s, 0.5, 6.4, 12.3, 1.0,
            "metaDEBASS today reports calibrated expert_confidence from 12 trust heads across 7 brokers and two surveys. "
            "The secondary follow-up proxy reaches AUC 0.970 with calibrated ECE 0.006. "
            "The remaining 16 registered classifiers are waiting on broker backfill, upstream release, or LSST scaling.",
            size=13, color=BLACK)


# ---------- brand mark / slide indicator on every non-title slide ----------
def add_brand_mark(slide, slide_num, total):
    add_textbox(slide, 0.5, 0.05, 8.0, 0.2,
                "metaDEBASS  ·  v7",
                size=8, color=GREY)
    add_textbox(slide, 8.5, 0.05, 4.3, 0.2,
                f"slide {slide_num} / {total}",
                size=8, color=GREY, align=PP_ALIGN.RIGHT)


total_slides = len(prs.slides)
for _i, _slide in enumerate(prs.slides, start=1):
    if _i == 1:
        continue
    add_brand_mark(_slide, _i, total_slides)


# ---------- save ----------
out_path = Path(__file__).resolve().parent / "metaDEBASS_status_2026-04-26.pptx"
prs.save(str(out_path))
print(f"Wrote {out_path}")
print(f"Slides: {len(prs.slides)}")
